# src/llm_parser.py
"""This script performs an advanced, LLM-based parsing of PDF and other document types.

It processes files in a source directory including PDFs, DOCX, PPTX, CSVs, images,
plain text, HTML, and EPUB, performs OCR (if needed),
and uses a large language model (LLM) to generate structured summaries.

The output is saved in two formats:
1. A detailed JSON file containing all results.
2. Individual Markdown files for each document, providing a human-readable
   summary with metadata and collapsible sections for the raw extracted text.
"""

import json
import sys
from datetime import datetime
from pathlib import Path

import cv2
import hydra
import numpy as np
import pandas as pd
import pytesseract
from bs4 import BeautifulSoup
from docx import Document
from ebooklib import epub
from langchain.prompts import PromptTemplate
from langchain_ollama import OllamaLLM
from omegaconf import DictConfig, OmegaConf
from pdf2image import convert_from_path
from pptx import Presentation
from PyPDF2 import PdfReader

from src.utils import RichConsoleManager


@hydra.main(version_base=None, config_path="../configs", config_name="rag_pipeline")
def parse_data_with_llms(config: DictConfig):
    console = RichConsoleManager.get_console()
    console.print(OmegaConf.to_yaml(config), style="warning")

    input_folder = Path(config.dataset_path)
    output_markdown_dir = Path(config.dataset_path)
    output_json_path = Path(config.llm_parser.output_json)
    output_markdown_dir.mkdir(exist_ok=True)

    if config.llm_parser.save_imgs:
        imgs_dir = Path(config.llm_parser.imgs_dir)
        imgs_dir.mkdir(exist_ok=True)
    else:
        imgs_dir = None

    llm = OllamaLLM(model=config.model.name)
    prompt = PromptTemplate.from_template(
        """
    You are a smart assistant tasked with reading diverse documents, such as:

    - Technical hardware datasheets
    - Sensor specifications
    - Instruction manuals
    - Books (educational, dietary, technical)
    - Game rulebooks
    - Research guidelines

    The source is a document titled: **{document_name}**,
    found in the folder: **{category}**

    The document may be written in English, Spanish, French,
    German, Italian, Portuguese, or another language.
    Please process the text accordingly and respond in the same language.

    Below is the content from page {page_number}:
    ---
    {page_text}
    ---

    Please summarize the key ideas or rules in **clear bullet points**.
    Adapt to the document type:
    - If it’s a **datasheet**, extract specs, features, and usage notes.
    - If it’s a **rulebook**, extract gameplay rules, setup, and conditions.
    - If it’s a **book**, extract core arguments, teachings, or steps.
    - If it’s a **guide/manual**, summarize the procedures or best practices.
    - If it’s a **scientific paper**, extract objective, methodology, and findings.

    Also try to extract metadata such as the author(s), publication date,
    or source if available.

    Avoid page numbers, headers, or formatting artifacts.
    """
    )
    runnable = prompt | llm

    summary_results = []
    supported_exts = [
        ".pdf",
        ".docx",
        ".pptx",
        ".csv",
        ".jpg",
        ".jpeg",
        ".png",
        ".txt",
        ".html",
        ".htm",
        ".epub",
    ]

    for path in input_folder.rglob("*"):
        if path.suffix.lower() not in supported_exts:
            continue

        category = path.parent.name
        md_filename = path.parent / f"{path.name}.md"
        if md_filename.exists():
            console.print(f"Skipping {path.name} - already processed.", style="info")
            continue

        console.print(f"\nProcessing: {category}/{path.name}", style="info")
        doc_summaries = []
        toc_entries = []
        metadata = {}

        texts = []
        ext = path.suffix.lower()

        if ext == ".pdf":
            try:
                reader = PdfReader(path)
                doc_info = reader.metadata
                if doc_info:
                    metadata = {
                        "author": doc_info.author,
                        "creator": doc_info.creator,
                        "producer": doc_info.producer,
                        "subject": doc_info.subject,
                        "title": doc_info.title,
                        "created": (
                            str(doc_info.creation_date)
                            if doc_info.creation_date
                            else None
                        ),
                    }
            except Exception as e:
                metadata = {"error": f"Metadata extraction failed: {e}"}

            images = convert_from_path(
                str(path),
                dpi=300,
                output_folder=str(imgs_dir) if imgs_dir else None,
                fmt="png",
            )
            for i, img in enumerate(images):
                img_cv = cv2.cvtColor(np.array(img), cv2.COLOR_RGB2BGR)
                try:
                    osd = pytesseract.image_to_osd(
                        img_cv, output_type=pytesseract.Output.DICT
                    )
                    rotation = osd["rotate"]
                    if rotation > 0:
                        (h, w) = img_cv.shape[:2]
                        M = cv2.getRotationMatrix2D((w // 2, h // 2), -rotation, 1.0)
                        img_cv = cv2.warpAffine(
                            img_cv, M, (w, h), flags=cv2.INTER_CUBIC
                        )
                except:
                    pass
                gray = cv2.cvtColor(img_cv, cv2.COLOR_BGR2GRAY)
                enhanced = cv2.adaptiveThreshold(
                    gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 35, 15
                )
                text = pytesseract.image_to_string(
                    enhanced, lang=config.llm_parser.tesseract_languages
                ).strip()
                if text:
                    texts.append((i + 1, text))

        elif ext in [".jpg", ".jpeg", ".png"]:
            img = cv2.imread(str(path))
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
            enhanced = cv2.adaptiveThreshold(
                gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 35, 15
            )
            text = pytesseract.image_to_string(
                enhanced, lang=config.llm_parser.tesseract_languages
            ).strip()
            if text:
                texts.append((1, text))

        elif ext == ".docx":
            doc = Document(path)
            text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            texts.append((1, text))

        elif ext == ".pptx":
            prs = Presentation(path)
            slide_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
            texts.append((1, "\n".join(slide_text)))

        elif ext == ".csv":
            df = pd.read_csv(path)
            texts.append((1, df.to_string(index=False)))

        elif ext == ".txt":
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                texts.append((1, f.read()))

        elif ext in [".html", ".htm"]:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
                texts.append((1, soup.get_text()))

        elif ext == ".epub":
            try:
                book = epub.read_epub(str(path))
                epub_texts = []
                for doc in book.get_items():
                    if doc.get_type() == epub.ITEM_DOCUMENT:
                        soup = BeautifulSoup(doc.get_body_content(), "html.parser")
                        epub_texts.append(soup.get_text())
                texts.append((1, "\n".join(epub_texts)))
            except Exception as e:
                texts.append((1, f"[EPUB parsing error: {str(e)}]"))

        for page_num, page_text in texts:
            summary = runnable.invoke(
                {
                    "document_name": path.name,
                    "category": category,
                    "page_number": page_num,
                    "page_text": page_text,
                }
            )
            result = {
                "document": path.name,
                "category": category,
                "page": page_num,
                "summary": summary,
                "ocr_text": page_text,
                "metadata": metadata,
            }
            summary_results.append(result)
            doc_summaries.append(result)
            toc_entries.append(f"- [Page {page_num}](#page-{page_num})")

        with open(md_filename, "w", encoding="utf-8") as f:
            f.write(f"# Summary for {path.name}\n")
            f.write(f"_Generated: {datetime.now().isoformat()}\n\n")
            if metadata:
                f.write("## Document Metadata\n")
                for k, v in metadata.items():
                    if v:
                        f.write(f"- **{k.capitalize()}**: {v}\n")
                f.write("\n")
            f.write("## Table of Contents\n")
            f.write("\n".join(toc_entries) + "\n\n")
            for item in doc_summaries:
                f.write(f"## Page {item['page']}\n\n")
                f.write(f"**Summary:**\n\n{item['summary']}\n\n")
                f.write(
                    "<details><summary>Raw OCR/Text</summary>\n\n"
                    f" ```{item['ocr_text']}```\n</details>\n\n"
                )

    with open(output_json_path, "w", encoding="utf-8") as f:
        json.dump(summary_results, f, ensure_ascii=False, indent=2)

    console.print(
        f"\n Markdown summaries saved in `{output_markdown_dir}`.\n"
        f" Full data saved in `{output_json_path}`."
    )


if __name__ == "__main__":
    research_dir = Path(__file__).resolve().parent
    sys.argv.append(f"hydra.run.dir={research_dir}")
    sys.argv.append("hydra.output_subdir=null")
    parse_data_with_llms()
