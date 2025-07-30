# rich_console_utils.py

import json
import logging
import pathlib
import sys
from datetime import datetime
from functools import lru_cache
from pathlib import Path
from typing import Dict, Optional, Union

import cv2
import numpy as np
import pandas as pd
import pytesseract
from bs4 import BeautifulSoup
from docx import Document
from ebooklib import epub
from omegaconf import DictConfig
from pdf2image import convert_from_path
from pptx import Presentation
from PyPDF2 import PdfReader
from rich.console import Console
from rich.logging import RichHandler
from rich.theme import Theme


class RichConsoleManager:
    """Provides a memoized and configurable Rich Console instance, with optional file
    logging and logging module integration."""

    DEFAULT_THEME = {
        "info": "bold bright_green",
        "warning": "bold bright_yellow",
        "danger": "bold bright_red",
        "summary": "italic green",
    }

    @classmethod
    def get_console(
        cls,
        theme_overrides: Optional[Dict[str, str]] = None,
        record: bool = False,
        log_path: Optional[Union[str, pathlib.Path]] = None,
    ) -> Console:
        """Get a Rich Console instance.

        Args:
            theme_overrides: Custom styles to override the defaults.
            record: Enable console recording (for testing/logging).
            log_path: If set, console output will be written to the given file.

        Returns:
            Console: Configured Rich Console instance.
        """
        theme = cls._build_theme(theme_overrides)
        log_path_str = str(log_path) if log_path else None
        return cls._get_cached_console(theme, record, log_path_str)

    @classmethod
    def setup_logging(
        cls,
        level: int = logging.INFO,
        console: Optional[Console] = None,
        log_format: str = "%(message)s",
        show_path: bool = False,
    ) -> None:
        """Setup logging with RichHandler.

        Args:
            level: Logging level (e.g., logging.DEBUG).
            console: Optional Rich Console to attach (default will be used otherwise).
            log_format: Log format string.
            show_path: Whether to display file path in log messages.
        """
        logging.basicConfig(
            level=level,
            format=log_format,
            datefmt="[%X]",
            handlers=[
                RichHandler(
                    console=console,
                    show_path=show_path,
                    markup=True,
                )
            ],
        )

    @staticmethod
    def _build_theme(overrides: Optional[Dict[str, str]]) -> Dict[str, str]:
        theme = RichConsoleManager.DEFAULT_THEME.copy()
        if overrides:
            theme.update(overrides)
        return theme

    @staticmethod
    @lru_cache(maxsize=None)
    def _cached_console_factory(
        theme_dict_frozen: frozenset,
        record: bool,
        log_path: Optional[str],
    ) -> Console:
        theme = Theme(dict(theme_dict_frozen))
        file = open(log_path, "a") if log_path else sys.stdout
        return Console(theme=theme, record=record, file=file)

    @classmethod
    def _get_cached_console(
        cls,
        theme_dict: Dict[str, str],
        record: bool,
        log_path: Optional[str],
    ) -> Console:
        frozen = frozenset(theme_dict.items())
        return cls._cached_console_factory(
            theme_dict_frozen=frozen, record=record, log_path=log_path
        )


def generate_markdown_from_json(json_path: Path, md_path: Path):
    """Generates a human-readable Markdown file from a structured JSON summary file.

    Args:
        json_path (Path): The path to the input JSON file.
        md_path (Path): The path to the output Markdown file to be created.
    """
    try:
        with open(json_path, "r", encoding="utf-8") as f:
            data = json.load(f)

        if not data:
            return

        # Extract metadata and document name from the first entry
        first_item = data[0]
        doc_name = first_item.get("document", "Unknown Document")
        metadata = first_item.get("source_metadata", {})

        with open(md_path, "w", encoding="utf-8") as f:
            f.write(f"# Summary for {doc_name}\n")
            f.write(f"_Generated: {datetime.now().isoformat()}\n\n")

            if metadata and not metadata.get("error"):
                f.write("## 📄 Document Metadata\n")
                for key, val in metadata.items():
                    if val:
                        f.write(f"- **{key.capitalize()}**: {val}\n")
                f.write("\n")

            toc_entries = [
                f"- [Page {item['page']}](#page-{item['page']})" for item in data
            ]
            f.write("## 📑 Table of Contents\n")
            f.write("\n".join(toc_entries) + "\n\n")

            for item in data:
                f.write("--- \n\n")
                f.write(f"### Page {item['page']}\n\n")

                # Get the summary, which we know is a string.
                summary_text = item.get("summary", "No summary was generated.")

                # Write the summary text directly.
                f.write(f"**Summary:**\n\n{summary_text}\n\n")

                # Write the collapsible raw text
                ocr_text = item.get("ocr_text", "")
                f.write(
                    f"<details><summary>Raw OCR/Text</summary>\n\n"
                    f" ```\n{ocr_text}\n```\n</details>\n\n"
                )

    except Exception as e:
        print(
            f"  - WARNING: Could not generate Markdown for {md_path.name}. Reason: {e}"
        )


def process_single_document(path: Path, config: DictConfig, console):
    """Orchestrates the text and metadata extraction for a single document.

    This function no longer calls the LLM. It focuses solely on robustly
    extracting the raw text from all pages of a given document.

    Returns:
        tuple: A tuple containing (list_of_page_texts, document_metadata),
               where list_of_page_texts is a list of (page_num, text) tuples.
    """
    console.print(
        f"\nExtracting text from: {path.relative_to(Path(config.dataset_path).parent)}",
        style="info",
    )

    metadata = {}
    texts = []  # This will be a list of (page_num, page_text) tuples
    ext = path.suffix.lower()

    if ext == ".pdf":
        try:
            reader = PdfReader(path)
            doc_info = reader.metadata
            if doc_info:
                metadata = {
                    "author": doc_info.author,
                    "creator": doc_info.creator,
                    "producer": doc_info.producer,
                    "subject": doc_info.subject,
                    "title": doc_info.title,
                    "created": (
                        str(doc_info.creation_date) if doc_info.creation_date else None
                    ),
                }
        except Exception as e:
            metadata = {"error": f"Metadata extraction failed: {e}"}

        images = convert_from_path(str(path), dpi=300, fmt="png")
        for i, img in enumerate(images):
            img_cv = cv2.cvtColor(np.array(img), cv2.COLOR_RGB2BGR)
            try:
                osd = pytesseract.image_to_osd(
                    img_cv, output_type=pytesseract.Output.DICT
                )
                rotation = osd["rotate"]
                if rotation > 0:
                    (h, w) = img_cv.shape[:2]
                    M = cv2.getRotationMatrix2D((w // 2, h // 2), -rotation, 1.0)
                    img_cv = cv2.warpAffine(img_cv, M, (w, h), flags=cv2.INTER_CUBIC)
            except:
                pass
            gray = cv2.cvtColor(img_cv, cv2.COLOR_BGR2GRAY)
            enhanced = cv2.adaptiveThreshold(
                gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 35, 15
            )
            tesseract_config = f"--psm 4 -l {config.llm_parser.tesseract_languages}"
            text = pytesseract.image_to_string(
                enhanced, config=tesseract_config
            ).strip()
            if text:
                texts.append((i + 1, text))

    elif ext in [".jpg", ".jpeg", ".png"]:
        img = cv2.imread(str(path))
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        enhanced = cv2.adaptiveThreshold(
            gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 35, 15
        )
        text = pytesseract.image_to_string(
            enhanced, lang=config.llm_parser.tesseract_languages
        ).strip()
        if text:
            texts.append((1, text))

    elif ext == ".docx":
        doc = Document(path)
        text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        texts.append((1, text))

    elif ext == ".pptx":
        prs = Presentation(path)
        slide_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
        texts.append((1, "\n".join(slide_text)))

    elif ext == ".csv":
        df = pd.read_csv(path)
        texts.append((1, df.to_string(index=False)))

    elif ext == ".txt":
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            texts.append((1, f.read()))

    elif ext in [".html", ".htm"]:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f.read(), "html.parser")
            texts.append((1, soup.get_text()))

    elif ext == ".epub":
        try:
            book = epub.read_epub(str(path))
            epub_texts = []
            for doc in book.get_items():
                if doc.get_type() == epub.ITEM_DOCUMENT:
                    soup = BeautifulSoup(doc.get_body_content(), "html.parser")
                    epub_texts.append(soup.get_text())
            texts.append((1, "\n".join(epub_texts)))
        except Exception as e:
            texts.append((1, f"[EPUB parsing error: {str(e)}]"))

    return texts, metadata
